import pptx
import pdfplumber


# from unstructured.partition.pdf import partition_pdf

def ppt_table_to_markdown(table):
    rows = table.rows
    num_rows = len(rows)
    num_columns = len(rows[0].cells)

    markdown_table = []
    first_line = True

    for row in rows:
        if first_line:
            header_row = [cell.text.strip() for cell in rows[0].cells]
            markdown_table.append('| ' + ' | '.join(header_row) + ' |')
            markdown_table.append('| ' + ' | '.join(['---'] * num_columns) + ' |')
            first_line = False
        row_data = [cell.text.strip() for cell in row.cells]
        markdown_table.append('| ' + ' | '.join(row_data) + ' |')

    return '\n'.join(markdown_table)


class read_file():
    def __init__(self, location):
        if location.endswith('.pptx') or location.endswith('.ppt'):
            self.mode = 'ppt'
            self.slides = pptx.Presentation(location).slides
        if location.endswith('.pdf'):
            self.mode = 'pdf'
            self.path = location

    def fit(self, trucation=0, verbose=False):
        if self.mode == 'ppt':
            content = ''
            start = 0

            for slide in self.slides:
                start += 1
                temp_content = ''
                for shape in slide.shapes:
                    if shape.has_text_frame:

                        text_frame = shape.text_frame

                        for paragraph in text_frame.paragraphs:
                            temp_content += paragraph.text + ' '

                    if shape.has_table:
                        table = shape.table
                        markdown_table = ppt_table_to_markdown(table)
                        temp_content += markdown_table
                    # print(markdown_table)
                    # print('---') 
                    if verbose:
                        print('temp', temp_content)
                if len(temp_content) > trucation:
                    content += f"第{start}页：\n"
                    content += temp_content
                    content += '\n'
                    if verbose:
                        print('---------------------')
            return content
        elif self.mode == 'pdf':
            with pdfplumber.open(self.path) as pdf:
                content = ''
                for i in range(len(pdf.pages)):
                    # 读取PDF文档第i+1页
                    page = pdf.pages[i]
                    # page.extract_text()函数即读取文本内容，下面这步是去掉文档最下面的页码
                    page_content = '\n'.join(page.extract_text().split('\n')[:-1])
                    content += f"\n第{i}页：\n"
                    content = content + page_content

            # test1 = partition_pdf(filename=self.path)
            # content = ''
            # pages = 0
            # for item in test1:
            #     if item.metadata.page_number != pages:
            #         content += f"第{pages+1}页：\n"
            #         pages += 1
            #         if verbose:
            #             print('---------------------')
            #             print(f"\n第{pages+1}页：\n")
            #     if len(item.text) > trucation:
            #         content+=item.text
            #         content+='\n'
            #         if verbose:
            #             print(item.text)
            return content


if __name__ == '__main__':
    location = input('location:')
    read_file = read_file(location)
    temp = read_file.fit(trucation=0, verbose=True)
    print(temp)
